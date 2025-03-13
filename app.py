import string
import nltk
import pdfplumber
import pandas as pd
import numpy as np
import streamlit as st
import matplotlib.pyplot as plt
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.ensemble import RandomForestClassifier
from sklearn.pipeline import make_pipeline
from nltk.corpus import stopwords
from streamlit import session_state as ss
from streamlit_pdf_viewer import pdf_viewer
from docx import Document
import pptx
from PIL import Image
import pytesseract
import io
import speech_recognition as sr
from pydub import AudioSegment
import tempfile
from sklearn.metrics.pairwise import cosine_similarity

nltk.download('stopwords')

# Expanded training data for better classification accuracy
DOCUMENT_TYPES = [
    "Resume", "Research Paper", "Scientific Report", "Medical Report", "Financial Report",
    "Legal Document", "Business Proposal", "Technical Manual", "Educational Material", "Audio Transcript"
]

training_texts = [
    "Experience, education, skills, projects, certifications",  # Resume
    "Abstract, introduction, methodology, results, conclusion",  # Research Paper
    "Experiment, observation, analysis, hypothesis",  # Scientific Report
    "Patient details, diagnosis, treatment, medication",  # Medical Report
    "Balance sheet, profit and loss, financial statement",  # Financial Report
    "Contract, agreement, clauses, legal terms",  # Legal Document
    "Market analysis, financial projection, business strategy",  # Business Proposal
    "Hardware specifications, installation guide, troubleshooting",  # Technical Manual
    "Syllabus, learning objectives, course material",  # Educational Material
    "Recorded speech, conversations, audio transcripts"  # Audio Transcript
]

classifier = make_pipeline(TfidfVectorizer(), RandomForestClassifier(n_estimators=300, random_state=42))
classifier.fit(training_texts, DOCUMENT_TYPES)

def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + " "
    return text.strip()

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file."""
    doc = Document(docx_file)
    return ' '.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    """Extract text from a PPTX file."""
    presentation = pptx.Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text.strip()

def extract_text_from_image(image_file):
    """Extract text from an image file."""
    image = Image.open(io.BytesIO(image_file.read()))
    return pytesseract.image_to_string(image)

def extract_text_from_audio(audio_file):
    """Extract text from an audio file."""
    recognizer = sr.Recognizer()
    audio_format = audio_file.type.split('/')[1]

    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{audio_format}') as temp_audio:
        temp_audio.write(audio_file.read())
        temp_audio_path = temp_audio.name

    if audio_format in ['mp3', 'mp4']:
        audio = AudioSegment.from_file(temp_audio_path)
        temp_audio_path = temp_audio_path.replace(audio_format, 'wav')
        audio.export(temp_audio_path, format='wav')

    with sr.AudioFile(temp_audio_path) as source:
        audio_data = recognizer.record(source)

    return recognizer.recognize_google(audio_data)

def preprocess_text(text):
    """Preprocess text by removing punctuation, stopwords, and converting to lowercase."""
    text = text.replace('-', ' ')
    text = text.translate(str.maketrans('', '', string.punctuation))
    text = text.lower()
    stop_words = set(stopwords.words('english'))
    words = text.split()
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

def tfidf_vectorization(text_data):
    """Compute TF-IDF scores and return a sorted DataFrame."""
    vectorizer = TfidfVectorizer(ngram_range=(1, 4), stop_words='english')
    tfidf_matrix = vectorizer.fit_transform(text_data)
    feature_names = vectorizer.get_feature_names_out()
    df_tfidf = pd.DataFrame(tfidf_matrix.toarray(), columns=feature_names)
    df_tfidf_sum = df_tfidf.sum(axis=0).reset_index()
    df_tfidf_sum.columns = ['Phrase', 'TF-IDF Score']
    df_tfidf_sorted = df_tfidf_sum.sort_values(by='TF-IDF Score', ascending=False).reset_index(drop=True)
    df_tfidf_sorted.index = df_tfidf_sorted.index + 1
    df_tfidf_sorted['Phrase'] = df_tfidf_sorted['Phrase'].str.title()
    return df_tfidf_sorted

def classify_document(text):
    """Classify the document type using the trained classifier."""
    return classifier.predict([text])[0]

def compute_similarity(text1, text2):
    """Compute cosine similarity between two texts."""
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform([text1, text2])
    return cosine_similarity(tfidf_matrix[0], tfidf_matrix[1])[0][0]

def main():
    """Streamlit app for Smart Document Analyzer."""
    st.title("Smart Document Analyzer")

    uploaded_files = st.sidebar.file_uploader("Drop your Files Here", type=['pdf', 'docx', 'pptx', 'png', 'jpeg', 'jpg', 'mp3', 'mp4', 'wav'], accept_multiple_files=True)

    if uploaded_files:
        texts = []

        for uploaded_file in uploaded_files:
            file_type = uploaded_file.type
            if 'pdf' in file_type:
                text = extract_text_from_pdf(uploaded_file)
            elif 'officedocument.wordprocessingml' in file_type:
                text = extract_text_from_docx(uploaded_file)
            elif 'officedocument.presentationml' in file_type:
                text = extract_text_from_pptx(uploaded_file)
            elif 'image' in file_type:
                text = extract_text_from_image(uploaded_file)
            elif 'audio' in file_type:
                text = extract_text_from_audio(uploaded_file)
            else:
                st.error("Unsupported file type")
                continue

            texts.append(preprocess_text(text))

        if len(texts) == 2:
            similarity_score = compute_similarity(texts[0], texts[1])
            st.subheader("🔗 Document Similarity Score")
            st.write(f"**{similarity_score:.2f}**")

if __name__ == "__main__":
    main()
